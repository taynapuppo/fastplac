import io
import os
import json
import time
import copy
import tempfile
import requests as http_requests
import google.auth.transport.requests
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseUpload
from pypdf import PdfWriter, PdfReader
from pptx import Presentation as PptxPresentation
from lxml import etree

SCOPES = [
    "https://www.googleapis.com/auth/drive",
    "https://www.googleapis.com/auth/presentations",
]

_BASE_DIR    = os.path.dirname(os.path.dirname(__file__))
_SECRETS_DIR = os.path.join(_BASE_DIR, "secrets")
TOKEN_FILE   = os.path.join(_SECRETS_DIR, "token.json")
CREDS_FILE   = os.path.join(_SECRETS_DIR, "client_secret.json")


# ─────────────────────────────────────────────────────────────
#  Autenticação
# ─────────────────────────────────────────────────────────────

def get_services():
    creds = None

    try:
        import streamlit as st
        if "google" in st.secrets and "token" in st.secrets["google"]:
            token_info = json.loads(st.secrets["google"]["token"])
            creds = Credentials.from_authorized_user_info(token_info, SCOPES)
    except Exception:
        pass

    if creds is None and os.path.exists(TOKEN_FILE):
        creds = Credentials.from_authorized_user_file(TOKEN_FILE, SCOPES)

    if creds and creds.expired and creds.refresh_token:
        creds.refresh(google.auth.transport.requests.Request())
        try:
            with open(TOKEN_FILE, "w") as f:
                f.write(creds.to_json())
        except Exception:
            pass

    if not creds or not creds.valid:
        try:
            import streamlit as st
            if "google" in st.secrets and "client_secret" in st.secrets["google"]:
                client_info = json.loads(st.secrets["google"]["client_secret"])
                with tempfile.NamedTemporaryFile(mode="w", suffix=".json", delete=False) as tmp:
                    json.dump(client_info, tmp)
                    tmp_path = tmp.name
                flow = InstalledAppFlow.from_client_secrets_file(tmp_path, SCOPES)
                os.unlink(tmp_path)
            else:
                flow = InstalledAppFlow.from_client_secrets_file(CREDS_FILE, SCOPES)
        except Exception:
            flow = InstalledAppFlow.from_client_secrets_file(CREDS_FILE, SCOPES)

        creds = flow.run_local_server(port=0)
        os.makedirs(_SECRETS_DIR, exist_ok=True)
        with open(TOKEN_FILE, "w") as f:
            f.write(creds.to_json())

    drive  = build("drive",  "v3", credentials=creds)
    slides = build("slides", "v1", credentials=creds)
    return drive, slides, creds


# ─────────────────────────────────────────────────────────────
#  Retry para Drive API (ainda usada para upload/export)
# ─────────────────────────────────────────────────────────────

def _execute_with_retry(request, max_retries: int = 7):
    for attempt in range(max_retries):
        try:
            return request.execute()
        except HttpError as e:
            if e.resp.status in (429, 500, 503) and attempt < max_retries - 1:
                wait = min(2 ** (attempt + 2), 60)
                print(f"[RETRY {e.resp.status}] Aguardando {wait}s (tentativa {attempt + 1}/{max_retries})...")
                time.sleep(wait)
            else:
                raise


# ─────────────────────────────────────────────────────────────
#  Download de template como PPTX (com cache por tipo)
# ─────────────────────────────────────────────────────────────

def _download_template_pptx(creds, template_id: str) -> bytes:
    """Baixa um Google Slides como PPTX. Usa autenticação OAuth2."""
    if creds.expired and creds.refresh_token:
        creds.refresh(google.auth.transport.requests.Request())
    url  = f"https://docs.google.com/presentation/d/{template_id}/export/pptx"
    resp = http_requests.get(
        url,
        headers={"Authorization": f"Bearer {creds.token}"},
        timeout=60,
    )
    resp.raise_for_status()
    return resp.content


# ─────────────────────────────────────────────────────────────
#  Preenchimento LOCAL de placeholders via python-pptx
# ─────────────────────────────────────────────────────────────

def _normalizar_runs(paragraph):
    """
    Consolida o texto de todos os runs de um parágrafo em um único run,
    preservando a formatação do primeiro run.
    Necessário porque o Google Slides frequentemente fragmenta texto
    como {{chave}} em múltiplos runs ao exportar como PPTX.
    """
    runs = paragraph.runs
    if len(runs) <= 1:
        return

    full_text = "".join(r.text or "" for r in runs)
    runs[0].text = full_text
    for r in runs[1:]:
        r.text = ""


def _iter_all_shapes(shapes):
    """
    Itera todos os shapes recursivamente, entrando em grupos (MSO_SHAPE_TYPE.GROUP = 6).
    A API do Slides fazia replaceAllText em toda a apresentação, incluindo shapes
    agrupados — este iterador replica esse comportamento.
    """
    for shape in shapes:
        yield shape
        if shape.shape_type == 6:          # GROUP
            yield from _iter_all_shapes(shape.shapes)


def _substituir_em_text_frame(text_frame, replacements: dict):
    """
    Aplica normalização e substituição em todos os parágrafos de um text_frame.
    Usa re.IGNORECASE para replicar o comportamento padrão da Slides API
    (matchCase: false) — ex: {{CLIENTE}} bate com a chave 'Cliente'.
    """
    import re
    for para in text_frame.paragraphs:
        _normalizar_runs(para)
        for run in para.runs:
            text = run.text or ""
            for key, value in replacements.items():
                pattern = re.escape(f"{{{{{key}}}}}")
                text = re.sub(pattern, lambda _: value, text, flags=re.IGNORECASE)
            run.text = text


def _fill_pptx_placeholders(pptx_bytes: bytes, data: dict) -> bytes:
    """
    Preenche placeholders no formato {{chave}} em um arquivo PPTX.
    Trabalha inteiramente em memória, sem nenhuma chamada de API.
    Itera recursivamente em group shapes, replicando o comportamento
    do replaceAllText da Slides API.
    """
    replacements = {
        k: (str(v).strip() if v not in (None, "") else "")
        for k, v in data.items()
    }

    prs = PptxPresentation(io.BytesIO(pptx_bytes))

    for slide in prs.slides:
        for shape in _iter_all_shapes(slide.shapes):
            # Caixas de texto
            if shape.has_text_frame:
                _substituir_em_text_frame(shape.text_frame, replacements)

            # Tabelas
            if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                for row in shape.table.rows:
                    for cell in row.cells:
                        _substituir_em_text_frame(cell.text_frame, replacements)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# ─────────────────────────────────────────────────────────────
#  Duplicação LOCAL do primeiro slide (sem API)
# ─────────────────────────────────────────────────────────────

def _duplicate_first_slide_pptx(pptx_bytes: bytes, extra_copies: int) -> bytes:
    """
    Duplica o primeiro slide N vezes dentro de um PPTX, operando
    diretamente no ZIP/XML sem chamar nenhuma API.
    """
    if extra_copies <= 0:
        return pptx_bytes

    import zipfile
    import re

    NS_PPTX   = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NS_REL    = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    NS_CT     = "http://schemas.openxmlformats.org/package/2006/content-types"
    REL_SLIDE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
    CT_SLIDE  = "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"

    src_buf = io.BytesIO(pptx_bytes)
    out_buf = io.BytesIO()

    with zipfile.ZipFile(src_buf, "r") as src_zip:
        names = src_zip.namelist()

        # Encontra e lê o primeiro slide
        slide_names = sorted(
            [n for n in names if re.match(r"ppt/slides/slide[0-9]+\.xml$", n)],
            key=lambda x: int(re.search(r"[0-9]+", x).group()),
        )
        if not slide_names:
            return pptx_bytes

        first_slide = slide_names[0]
        first_slide_num = int(re.search(r"[0-9]+", first_slide).group())
        max_slide_num   = max(int(re.search(r"[0-9]+", n).group()) for n in slide_names)

        # Lê os XMLs que precisam ser modificados
        pres_root = etree.fromstring(src_zip.read("ppt/presentation.xml"))
        rels_root = etree.fromstring(src_zip.read("ppt/_rels/presentation.xml.rels"))
        ct_root   = etree.fromstring(src_zip.read("[Content_Types].xml"))

        sldIdLst = pres_root.find(f"{{{NS_PPTX}}}sldIdLst")
        max_id   = max(
            (int(e.get("id", 256)) for e in sldIdLst.findall(f"{{{NS_PPTX}}}sldId")),
            default=256,
        )
        max_rel_id = max(
            (int(re.sub(r"\D", "", e.get("Id", "0")) or 0) for e in rels_root),
            default=100,
        )

        extra_files = {}
        for i in range(extra_copies):
            new_num      = max_slide_num + i + 1
            new_path     = f"ppt/slides/slide{new_num}.xml"
            new_rel_path = f"ppt/slides/_rels/slide{new_num}.xml.rels"
            rel_id       = f"rId{max_rel_id + i + 1}"

            # Copia XML do slide
            extra_files[new_path] = src_zip.read(first_slide)

            # Copia .rels do slide (se existir)
            first_rel = f"ppt/slides/_rels/slide{first_slide_num}.xml.rels"
            if first_rel in names:
                extra_files[new_rel_path] = src_zip.read(first_rel)

            # Adiciona à lista de slides na apresentação
            max_id += 1
            el = etree.SubElement(sldIdLst, f"{{{NS_PPTX}}}sldId")
            el.set("id", str(max_id))
            el.set(f"{{{NS_REL}}}id", rel_id)

            # Adiciona relacionamento
            rel_el = etree.SubElement(rels_root, "Relationship")
            rel_el.set("Id", rel_id)
            rel_el.set("Type", REL_SLIDE)
            rel_el.set("Target", f"slides/slide{new_num}.xml")

            # Adiciona ao Content_Types
            ov = etree.SubElement(ct_root, f"{{{NS_CT}}}Override")
            ov.set("PartName", f"/{new_path}")
            ov.set("ContentType", CT_SLIDE)

        with zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED) as out_zip:
            for name in names:
                if name == "ppt/presentation.xml":
                    out_zip.writestr(
                        name,
                        etree.tostring(pres_root, xml_declaration=True, encoding="UTF-8", standalone=True),
                    )
                elif name == "ppt/_rels/presentation.xml.rels":
                    out_zip.writestr(
                        name,
                        etree.tostring(rels_root, xml_declaration=True, encoding="UTF-8", standalone=True),
                    )
                elif name == "[Content_Types].xml":
                    out_zip.writestr(
                        name,
                        etree.tostring(ct_root, xml_declaration=True, encoding="UTF-8", standalone=True),
                    )
                else:
                    out_zip.writestr(name, src_zip.read(name))

            for name, data in extra_files.items():
                out_zip.writestr(name, data)

    out_buf.seek(0)
    return out_buf.read()


# ─────────────────────────────────────────────────────────────
#  Mesclagem LOCAL de múltiplos PPTX
# ─────────────────────────────────────────────────────────────

def _merge_pptx(pptx_bytes_list: list[bytes]) -> bytes:
    """
    Mescla múltiplos PPTX trabalhando direto no ZIP/XML.
    Preserva slides, mídias e relacionamentos de cada arquivo.
    """
    import zipfile as _zf
    import re as _re

    NS_PPTX   = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NS_REL    = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    NS_CT     = "http://schemas.openxmlformats.org/package/2006/content-types"
    REL_SLIDE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
    CT_SLIDE  = "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"

    zips       = [_zf.ZipFile(io.BytesIO(b), "r") for b in pptx_bytes_list]
    base_zip   = zips[0]
    base_names = set(base_zip.namelist())
    out_buf    = io.BytesIO()

    slide_count = len([n for n in base_names if _re.match(r"ppt/slides/slide[0-9]+\.xml$", n)])
    media_names = {n for n in base_names if n.startswith("ppt/media/")}
    extra_files = {}
    new_pairs   = []  # (zip_path, rel_id)

    for src_zip in zips[1:]:
        src_names  = set(src_zip.namelist())
        src_slides = sorted(
            [n for n in src_names if _re.match(r"ppt/slides/slide[0-9]+\.xml$", n)],
            key=lambda x: int(_re.search(r"[0-9]+", x).group()),
        )
        for name in src_names:
            if name.startswith("ppt/media/") and name not in media_names:
                extra_files[name] = src_zip.read(name)
                media_names.add(name)
        for slide_path in src_slides:
            slide_count += 1
            new_path = f"ppt/slides/slide{slide_count}.xml"
            extra_files[new_path] = src_zip.read(slide_path)
            rel_src = slide_path.replace("ppt/slides/", "ppt/slides/_rels/") + ".rels"
            rel_dst = new_path.replace("ppt/slides/", "ppt/slides/_rels/") + ".rels"
            if rel_src in src_names:
                extra_files[rel_dst] = src_zip.read(rel_src)
            new_pairs.append((new_path, f"rId{100 + slide_count}"))

    pres_root = etree.fromstring(base_zip.read("ppt/presentation.xml"))
    sldIdLst  = pres_root.find(f"{{{NS_PPTX}}}sldIdLst")
    max_id    = max(
        (int(e.get("id", 256)) for e in sldIdLst.findall(f"{{{NS_PPTX}}}sldId")),
        default=256,
    )
    for _, rel_id in new_pairs:
        max_id += 1
        el = etree.SubElement(sldIdLst, f"{{{NS_PPTX}}}sldId")
        el.set("id", str(max_id))
        el.set(f"{{{NS_REL}}}id", rel_id)

    rels_root = etree.fromstring(base_zip.read("ppt/_rels/presentation.xml.rels"))
    for slide_path, rel_id in new_pairs:
        el = etree.SubElement(rels_root, "Relationship")
        el.set("Id", rel_id)
        el.set("Type", REL_SLIDE)
        el.set("Target", slide_path.replace("ppt/", ""))

    ct_root = etree.fromstring(base_zip.read("[Content_Types].xml"))
    for slide_path, _ in new_pairs:
        ov = etree.SubElement(ct_root, f"{{{NS_CT}}}Override")
        ov.set("PartName", f"/{slide_path}")
        ov.set("ContentType", CT_SLIDE)

    with _zf.ZipFile(out_buf, "w", _zf.ZIP_DEFLATED) as out_zip:
        for name in base_zip.namelist():
            if name == "ppt/presentation.xml":
                out_zip.writestr(name, etree.tostring(pres_root, xml_declaration=True, encoding="UTF-8", standalone=True))
            elif name == "ppt/_rels/presentation.xml.rels":
                out_zip.writestr(name, etree.tostring(rels_root, xml_declaration=True, encoding="UTF-8", standalone=True))
            elif name == "[Content_Types].xml":
                out_zip.writestr(name, etree.tostring(ct_root, xml_declaration=True, encoding="UTF-8", standalone=True))
            else:
                out_zip.writestr(name, base_zip.read(name))
        for name, data in extra_files.items():
            out_zip.writestr(name, data)

    for z in zips:
        z.close()
    out_buf.seek(0)
    return out_buf.read()


# ─────────────────────────────────────────────────────────────
#  Drive helpers
# ─────────────────────────────────────────────────────────────

def rename_file(drive, file_id: str, new_name: str) -> str:
    result = _execute_with_retry(
        drive.files().update(
            fileId=file_id,
            body={"name": new_name},
            fields="id, webViewLink",
        )
    )
    return result.get("webViewLink", "")


def delete_file(drive, file_id: str):
    try:
        drive.files().delete(fileId=file_id).execute()
    except HttpError as e:
        print(f"[WARN] Não foi possível deletar {file_id}: {e}")


def upload_pdf(drive, pdf_bytes: bytes, nome: str, folder_id: str) -> str:
    metadata = {"name": f"{nome}.pdf", "parents": [folder_id], "mimeType": "application/pdf"}
    media    = MediaIoBaseUpload(io.BytesIO(pdf_bytes), mimetype="application/pdf")
    arquivo  = _execute_with_retry(
        drive.files().create(body=metadata, media_body=media, fields="id, webViewLink")
    )
    return arquivo.get("webViewLink", "")


def export_as_pdf(creds, presentation_id: str) -> bytes:
    if creds.expired and creds.refresh_token:
        creds.refresh(google.auth.transport.requests.Request())
    url  = f"https://docs.google.com/presentation/d/{presentation_id}/export/pdf"
    resp = http_requests.get(
        url,
        headers={"Authorization": f"Bearer {creds.token}"},
        timeout=120,
    )
    resp.raise_for_status()
    return resp.content


def merge_pdfs(pdf_bytes_list: list[bytes]) -> bytes:
    writer = PdfWriter()
    for pdf_bytes in pdf_bytes_list:
        reader = PdfReader(io.BytesIO(pdf_bytes))
        for page in reader.pages:
            writer.add_page(page)
    out = io.BytesIO()
    writer.write(out)
    return out.getvalue()


# ─────────────────────────────────────────────────────────────
#  Função principal
# ─────────────────────────────────────────────────────────────

def gerar_pdf_consolidado(
    placas: list[dict],
    folder_id: str,
    template_ids: dict,
    nome_arquivo: str = "Placas",
    progress_callback=None,
) -> tuple[bytes, str, list[dict], str]:
    """
    Nova abordagem — sem batchUpdate, sem limite de quota:

    1. Baixa cada template UMA vez por tipo (cache em memória)
    2. Preenche placeholders LOCALMENTE via python-pptx
    3. Duplica slides localmente se Qtd > 1
    4. Mescla todos os PPTX localmente
    5. Faz 1 upload do PPTX mesclado → exporta como PDF
    6. Mantém o PPTX mesclado como Slides consolidado no Drive
    7. Faz 1 upload do PDF final

    Retorna: (pdf_bytes, link_pdf, slides_info, link_slides)
    """
    drive, _, creds = get_services()

    total = len(placas)

    # ── 1. Download dos templates (por tipo, com cache) ──
    if progress_callback:
        progress_callback(0.02, "Baixando templates...")

    tipos_unicos    = list({p["tipo"] for p in placas})
    template_cache: dict[str, bytes] = {}

    for i, tipo in enumerate(tipos_unicos):
        if progress_callback:
            progress_callback(
                0.02 + (i / len(tipos_unicos)) * 0.15,
                f"Baixando template {i + 1}/{len(tipos_unicos)}: {tipo}",
            )
        template_cache[tipo] = _download_template_pptx(creds, template_ids[tipo])
        time.sleep(0.3)  # gentil com a API de export

    # ── 2 + 3. Preenche e duplica cada placa localmente ──
    filled_pptx_list: list[bytes] = []
    slides_info:      list[dict]  = []

    for idx, placa in enumerate(placas):
        tipo  = placa["tipo"]
        dados = placa["dados"]
        qtd   = max(1, int(dados.get("Quantidade de Placas") or 1))

        if progress_callback:
            pct = 0.17 + (idx / total) * 0.55
            progress_callback(pct, f"Preenchendo placa {idx + 1}/{total}: {tipo}")

        # Preenche localmente (zero API calls)
        filled = _fill_pptx_placeholders(template_cache[tipo], dados)

        # Duplica slides localmente se necessário
        if qtd > 1:
            filled = _duplicate_first_slide_pptx(filled, qtd - 1)

        filled_pptx_list.append(filled)

        cliente = dados.get("Cliente", "")
        slides_info.append({"tipo": tipo, "cliente": cliente, "link": ""})

    # ── 4. Mescla todos os PPTX localmente ──
    if progress_callback:
        progress_callback(0.73, f"Mesclando {total} apresentações...")

    merged_pptx = _merge_pptx(filled_pptx_list)

    # ── 5. Faz upload do PPTX mesclado como Google Slides ──
    if progress_callback:
        progress_callback(0.80, "Enviando para o Google Drive...")

    nome_slides = f"Slides - {nome_arquivo}"
    metadata = {
        "name":     nome_slides,
        "parents":  [folder_id],
        "mimeType": "application/vnd.google-apps.presentation",
    }
    media = MediaIoBaseUpload(
        io.BytesIO(merged_pptx),
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        resumable=True,
    )
    result = _execute_with_retry(
        drive.files().create(body=metadata, media_body=media, fields="id, webViewLink")
    )
    slides_id   = result["id"]
    link_slides = result.get("webViewLink", "")

    # ── 6. Exporta como PDF a partir do Slides consolidado ──
    if progress_callback:
        progress_callback(0.90, "Exportando PDF...")

    # Aguarda um momento para o Drive processar a conversão
    time.sleep(3)
    pdf_bytes = export_as_pdf(creds, slides_id)

    # ── 7. Faz upload do PDF ──
    if progress_callback:
        progress_callback(0.96, "Salvando PDF na pasta de concluídos...")

    link_pdf = upload_pdf(drive, pdf_bytes, nome_arquivo, folder_id)

    if progress_callback:
        progress_callback(1.0, "Concluído!")

    return pdf_bytes, link_pdf, slides_info, link_slides